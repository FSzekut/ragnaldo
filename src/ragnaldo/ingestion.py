"""Pipeline offline de documentos, chunks, embeddings e índice FAISS."""

from __future__ import annotations

import hashlib
import json
import re
import warnings
from collections.abc import Callable, Iterable
from datetime import datetime, timezone
from pathlib import Path

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from ragnaldo.config import SETTINGS, VECTOR_STORE_DIR, RagSettings

MANIFEST_NAME = "artifact_manifest.json"

# Tabelas e planilhas viram muitos registros curtos. Agrupá-los em blocos evita
# um chunk por linha, que recupera mal e polui as citações.
CSV_ROWS_PER_BLOCK = 40
JSON_ITEMS_PER_BLOCK = 20


def _clean(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def load_pdf(path: Path) -> list[Document]:
    documents = PyPDFLoader(str(path)).load()
    for document in documents:
        page = document.metadata.get("page")
        if isinstance(page, int):
            document.metadata["location"] = f"página {page + 1}"
    return documents


def load_text(path: Path) -> list[Document]:
    return TextLoader(str(path), encoding="utf-8").load()


def load_html(path: Path) -> list[Document]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8"), "lxml")
    # Menu, rodapé e scripts se repetem em toda página e não carregam conteúdo:
    # como chunk, só competem com o texto real na busca por similaridade.
    for tag in soup(["script", "style", "nav", "footer", "header", "noscript", "form"]):
        tag.decompose()
    return [Document(page_content=_clean(soup.get_text("\n", strip=True)))]


def load_docx(path: Path) -> list[Document]:
    import docx

    document = docx.Document(str(path))
    blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))
    return [Document(page_content=_clean("\n".join(blocks)))]


def load_xlsx(path: Path) -> list[Document]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        documents = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value).strip() for value in row if value is not None]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                documents.append(
                    Document(
                        page_content=_clean("\n".join(rows)),
                        metadata={"location": f"planilha {sheet.title}"},
                    )
                )
        return documents
    finally:
        workbook.close()


def load_pptx(path: Path) -> list[Document]:
    from pptx import Presentation

    documents = []
    for number, slide in enumerate(Presentation(str(path)).slides, start=1):
        parts = [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if parts:
            documents.append(
                Document(
                    page_content=_clean("\n".join(parts)),
                    metadata={"location": f"slide {number}"},
                )
            )
    return documents


def load_csv(path: Path) -> list[Document]:
    import csv

    with path.open(encoding="utf-8", newline="") as handle:
        rows = [
            "; ".join(
                f"{column}: {value}"
                for column, value in record.items()
                if column and value and str(value).strip()
            )
            for record in csv.DictReader(handle)
        ]

    documents = []
    for start in range(0, len(rows), CSV_ROWS_PER_BLOCK):
        block = [row for row in rows[start : start + CSV_ROWS_PER_BLOCK] if row]
        if block:
            documents.append(
                Document(
                    page_content="\n".join(block),
                    metadata={"location": f"linhas {start + 1}-{start + len(block)}"},
                )
            )
    return documents


def load_json(path: Path) -> list[Document]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, list):
        return [Document(page_content=json.dumps(payload, ensure_ascii=False, indent=2))]

    documents = []
    for start in range(0, len(payload), JSON_ITEMS_PER_BLOCK):
        block = payload[start : start + JSON_ITEMS_PER_BLOCK]
        documents.append(
            Document(
                page_content=json.dumps(block, ensure_ascii=False, indent=2),
                metadata={"location": f"itens {start + 1}-{start + len(block)}"},
            )
        )
    return documents


# Fonte única da verdade sobre formatos: quem descobre e quem carrega leem esta
# mesma tabela, então é impossível um arquivo ser aceito e depois não ter leitor.
LOADERS: dict[str, Callable[[Path], list[Document]]] = {
    ".pdf": load_pdf,
    ".md": load_text,
    ".txt": load_text,
    ".html": load_html,
    ".htm": load_html,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
    ".csv": load_csv,
    ".json": load_json,
}


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file:
        for block in iter(lambda: file.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def discover_sources(*directories: Path) -> list[Path]:
    supported: list[Path] = []
    ignored: list[Path] = []

    for directory in directories:
        if not directory.exists():
            continue
        for path in directory.rglob("*"):
            if not path.is_file() or path.name.startswith("."):
                continue
            if path.suffix.lower() in LOADERS:
                supported.append(path)
            else:
                ignored.append(path)

    # Um arquivo baixado e nunca indexado é a falha mais cara aqui: não quebra
    # nada, some do corpus em silêncio e só aparece quando o agente não sabe
    # responder algo que deveria saber.
    if ignored:
        warnings.warn(
            "Arquivos sem loader registrado foram ignorados: "
            + ", ".join(sorted(path.name for path in ignored)),
            stacklevel=2,
        )

    return sorted(supported)


def load_sources(paths: Iterable[Path]) -> list[Document]:
    documents: list[Document] = []
    for path in paths:
        suffix = path.suffix.lower()
        loader = LOADERS.get(suffix)
        if loader is None:
            raise ValueError(f"Formato sem loader registrado: {path.name}")

        loaded = loader(path)
        source_hash = sha256_file(path)
        for document in loaded:
            document.metadata.setdefault("location", "documento")
            document.metadata.update(
                {
                    "source": path.name,
                    "source_path": str(path),
                    "source_sha256": source_hash,
                    "format": suffix.lstrip("."),
                }
            )
        documents.extend(loaded)
    return documents


def split_documents(
    documents: list[Document], settings: RagSettings = SETTINGS
) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        add_start_index=True,
    )
    chunks = splitter.split_documents(documents)
    for position, chunk in enumerate(chunks):
        identity = "|".join(
            [
                chunk.metadata.get("source_sha256", ""),
                str(chunk.metadata.get("location", "")),
                str(chunk.metadata.get("start_index", "")),
                chunk.page_content,
            ]
        )
        chunk.metadata["chunk_id"] = hashlib.sha256(identity.encode()).hexdigest()[:16]
        chunk.metadata["chunk_position"] = position
    return chunks


def create_embeddings(settings: RagSettings = SETTINGS) -> HuggingFaceEmbeddings:
    return HuggingFaceEmbeddings(
        model_name=settings.embedding_model,
        model_kwargs={"device": settings.device},
        encode_kwargs={"normalize_embeddings": True, "batch_size": 16},
    )


def build_index(
    chunks: list[Document],
    destination: Path = VECTOR_STORE_DIR,
    settings: RagSettings = SETTINGS,
) -> dict:
    if not chunks:
        raise ValueError("Nenhum chunk foi produzido; o índice não será criado.")

    destination.mkdir(parents=True, exist_ok=True)
    vector_store = FAISS.from_documents(chunks, create_embeddings(settings))
    vector_store.save_local(str(destination))

    artifact_hashes = {
        path.name: sha256_file(path)
        for path in destination.iterdir()
        if path.is_file() and path.name != MANIFEST_NAME
    }
    source_hashes = {
        chunk.metadata["source"]: chunk.metadata["source_sha256"] for chunk in chunks
    }
    manifest = {
        "created_at": datetime.now(timezone.utc).isoformat(),
        "embedding_model": settings.embedding_model,
        "chunk_size": settings.chunk_size,
        "chunk_overlap": settings.chunk_overlap,
        "chunk_count": len(chunks),
        "source_hashes": source_hashes,
        "artifact_hashes": artifact_hashes,
    }
    (destination / MANIFEST_NAME).write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return manifest


def verify_index(destination: Path = VECTOR_STORE_DIR) -> dict:
    manifest_path = destination / MANIFEST_NAME
    if not manifest_path.exists():
        raise FileNotFoundError(
            "Índice ainda não criado. Execute o notebook 01_ingestao_e_embeddings."
        )

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    for filename, expected_hash in manifest["artifact_hashes"].items():
        artifact = destination / filename
        if not artifact.is_file() or sha256_file(artifact) != expected_hash:
            raise RuntimeError(f"Artefato ausente ou alterado: {filename}")
    return manifest


def load_verified_index(
    destination: Path = VECTOR_STORE_DIR,
    settings: RagSettings = SETTINGS,
) -> tuple[FAISS, dict]:
    manifest = verify_index(destination)
    if manifest["embedding_model"] != settings.embedding_model:
        raise RuntimeError(
            "O modelo configurado difere daquele usado para construir o índice."
        )

    # FAISS.save_local usa pickle para o docstore. A desserialização só ocorre
    # depois que todos os artefatos são verificados contra o manifesto local.
    vector_store = FAISS.load_local(
        str(destination),
        create_embeddings(settings),
        allow_dangerous_deserialization=True,
    )
    return vector_store, manifest
